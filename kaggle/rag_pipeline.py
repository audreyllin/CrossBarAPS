# Standard library imports
import os
import json
import base64
import subprocess
import logging
import re
from constants import INDEX_VERSION_FILE, VALID_RATIOS, VALID_STYLES, ALLOWED_EXTENSIONS
from datetime import datetime, timezone
from pathlib import Path
import fcntl

# Third-party imports
import numpy as np
import requests
from filelock import FileLock

# Document processing imports
import pytesseract
import pdfplumber
import fitz
from PIL import Image
from docx import Document as DocxReader
from docx import Document as DocxWriter
from pptx import Presentation

# AI/Vector imports
import faiss
import replicate
from replicate.exceptions import ReplicateError
from openai import OpenAI

import zipfile
import tempfile
from pathlib import Path

ALLOWED_EXTENSIONS = {
    ".txt",
    ".pdf",
    ".doc",
    ".docx",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
    ".csv",
    ".tsv",
    ".jpg",
    ".jpeg",
    ".png",
    ".odt",
    ".rtf",
    ".zip",
}

# Constants and Configuration
INDEX_VERSION_FILE = "index_version.json"
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 300
EMBEDDING_MODEL = "text-embedding-3-large"
GPT_MODEL = "gpt-4o"
MIN_CHUNK_LENGTH = 100
MAX_FILE_SIZE_MB = 10
REPLICATE_API_TOKEN = os.environ.get("REPLICATE_API_TOKEN")
POSTER_MODEL = os.getenv("POSTER_MODEL", "ideogram-ai/ideogram-v3-turbo")
REPLICATE_USE_FILE_OUTPUT = False

# Initialize clients
client_oa = OpenAI()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("optimized_rag")

# ======================
# File Processing Functions
# ======================

def is_image_heavy_pdf(filepath):
    """Check if PDF is primarily image-based"""
    try:
        doc = fitz.open(filepath)
        image_page_count = sum(1 for p in doc if len(p.get_images(full=True)) > 0)
        return image_page_count / len(doc) > 0.5
    except:
        return False

def clean_text(text):
    """Clean and normalize text by removing extra whitespace"""
    return re.sub(r"\s+", " ", text).strip()

def extract_pdf_text(filepath):
    """Extract text from PDF, falling back to OCR if needed"""
    text = ""
    try:
        # First try regular text extraction
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""

        # Fallback to OCR if needed
        if not text.strip() or is_image_heavy_pdf(filepath):
            logger.info(f"Using OCR for PDF: {filepath}")
            doc = fitz.open(filepath)
            for page in doc:
                text += page.get_text() or ""
    except Exception as e:
        logger.error(f"PDF extraction failed: {str(e)}")
    return clean_text(text)[:50000]  # Limit to 50k characters


def extract_docx_text(filepath):
    """Extract text from Word documents"""
    try:
        doc = DocxReader(filepath)
        return clean_text("\n".join(p.text for p in doc.paragraphs))
    except Exception as e:
        logger.error(f"DOCX extraction failed: {str(e)}")
        return ""


def extract_text(filepath):
    """Main text extraction function that handles different file types"""
    ext = Path(filepath).suffix.lower()

    # Skip large files
    if os.path.getsize(filepath) > MAX_FILE_SIZE_MB * 1024 * 1024:
        logger.warning(f"Skipping large file: {filepath}")
        return ""

    try:
        if ext == ".pdf":
            return extract_pdf_text(filepath)
        elif ext in [".doc", ".docx"]:
            return extract_docx_text(filepath)
        elif ext in [".jpg", ".jpeg", ".png"]:
            return pytesseract.image_to_string(Image.open(filepath))
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return clean_text(f.read())[:50000]
        else:
            logger.warning(f"Unsupported file format: {ext}")
            return ""
    except Exception as e:
        logger.error(f"Extraction error for {filepath}: {str(e)}")
        return ""


def chunk_text(text):
    """Split text into chunks with overlap, respecting sentence boundaries"""
    if len(text) <= CHUNK_SIZE:
        return [text]

    chunks = []
    start = 0

    while start < len(text):
        end = min(start + CHUNK_SIZE, len(text))

        # Adjust to nearest sentence boundary if possible
        if end < len(text):
            boundary = max(
                text.rfind(".", start, end),
                text.rfind("?", start, end),
                text.rfind("!", start, end),
            )
            if boundary > start + MIN_CHUNK_LENGTH:
                end = boundary + 1

        chunks.append(text[start:end])
        start = end - CHUNK_OVERLAP

    return chunks


# ======================
# Index Management Functions
# ======================


def get_index_version():
    """Get current index version with file locking"""
    lock_path = os.path.join(OUTPUT_FOLDER, "index_version.lock")
    with FileLock(lock_path):
        try:
            with open(INDEX_VERSION_FILE, "r") as f:
                return json.load(f).get("version", 1)
        except:
            return 1


def increment_index_version():
    """Increment index version with file locking"""
    lock_path = os.path.join(OUTPUT_FOLDER, "index_version.lock")
    with FileLock(lock_path):
        try:
            with open(INDEX_VERSION_FILE, "r+") as f:
                version_data = json.load(f)
                version_data["version"] += 1
                f.seek(0)
                json.dump(version_data, f)
                f.truncate()
            return version_data["version"]
        except Exception as e:
            logger.error(f"Error incrementing index version: {str(e)}")
            return None


def get_embeddings(texts, api_key):
    """Generate embeddings using OpenAI's API"""
    client = OpenAI(api_key=api_key)
    try:
        response = client.embeddings.create(
            model=EMBEDDING_MODEL, input=texts, encoding_format="float"
        )
        return [data.embedding for data in response.data]
    except Exception as e:
        logger.error(f"Embedding generation failed: {str(e)}")
        return []


def create_index(api_key):
    """Create a FAISS index with version checking"""
    index_path = Path(OUTPUT_FOLDER) / "index.faiss"
    metadata_path = Path(OUTPUT_FOLDER) / "metadata.json"
    lock_path = Path(OUTPUT_FOLDER) / "index.lock"

    # Check if we need to rebuild
    current_version = get_index_version()

    with FileLock(str(lock_path)):
        if index_path.exists() and metadata_path.exists():
            # Verify index is up-to-date
            try:
                index = faiss.read_index(str(index_path))
                with open(metadata_path, "r") as f:
                    metadata = json.load(f)

                # Check if any files have changed since last index
                needs_rebuild = False
                for filepath in Path(UPLOAD_FOLDER).glob("*"):
                    if filepath.suffix.lower() in [".pdf", ".doc", ".docx", ".txt"]:
                        file_mtime = filepath.stat().st_mtime
                        if file_mtime > current_version:
                            needs_rebuild = True
                            break

                if not needs_rebuild:
                    return index, metadata
            except Exception as e:
                logger.warning(f"Error loading existing index, will rebuild: {str(e)}")

        # Rebuild index
        logger.info("Building new index")
        index = faiss.IndexFlatL2(3072)
        metadata = []

        for filepath in Path(UPLOAD_FOLDER).glob("*"):
            if filepath.suffix.lower() not in [".pdf", ".doc", ".docx", ".txt"]:
                continue

            logger.info(f"Processing: {filepath.name}")
            text = extract_text(str(filepath))
            if not text:
                continue

            chunks = chunk_text(text)
            embeddings = get_embeddings(chunks, api_key)

            for i, emb in enumerate(embeddings):
                if len(chunks[i]) > MIN_CHUNK_LENGTH:
                    index.add(np.array([emb]).astype("float32"))
                    metadata.append(
                        {
                            "filename": filepath.name,
                            "chunk_index": i,
                            "content": chunks[i][:500],
                            "version": current_version,
                        }
                    )

        # Save new index
        faiss.write_index(index, str(index_path))
        with open(metadata_path, "w") as f:
            json.dump(metadata, f)

        # Update version
        increment_index_version()

        return index, metadata


def search_index(index, metadata, query_embedding, k=5):
    """Search the FAISS index for relevant chunks"""
    distances, indices = index.search(np.array([query_embedding]).astype("float32"), k)
    return [metadata[i] for i in indices[0] if i < len(metadata)]


# ======================
# Media Generation Functions
# ======================


def _run_replicate(model_ref: str, input_data: dict) -> str:
    """Run a Replicate model and handle different output formats."""
    try:
        output = replicate.run(
            model_ref,
            input=input_data,
            use_file_output=REPLICATE_USE_FILE_OUTPUT,
        )
    except Exception as e:
        logger.error(f"Replicate run failed: {str(e)}")
        raise

    # Handle different output types
    if isinstance(output, list) and len(output) > 0:
        first = output[0]
    else:
        first = output

    # Duck-typing for URL extraction
    if hasattr(first, "url"):  # FileOutput object
        return first.url
    elif isinstance(first, str):
        return first
    elif isinstance(first, dict):
        # Look for URL in dictionary values
        for value in first.values():
            if isinstance(value, str) and value.startswith("http"):
                return value
            if hasattr(value, "url"):
                return value.url
        # Fallback to first value if no URL found
        first_val = next(iter(first.values()), None)
        return str(first_val) if first_val is not None else ""
    return str(output)  # Final fallback


def _run_dalle(prompt: str) -> str:
    """Fallback to DALL-E 3 if Replicate fails"""
    rsp = client_oa.images.generate(
        model="dall-e-3", prompt=prompt, size="1024x1024", n=1
    )
    return rsp.data[0].url


def _download(url: str, session_id: str, media_type: str) -> str:
    """Download remote media to /output and return local path."""
    safe_session = session_id or "anon"
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d%H%M%S")
    ext = ".mp4" if media_type == "video" else ".png"
    fn = f"{media_type}_{safe_session}_{timestamp}{ext}"
    out_path = os.path.join(OUTPUT_FOLDER, fn)

    with requests.get(url, stream=True, timeout=30) as r:
        r.raise_for_status()
        with open(out_path, "wb") as f:
            for chunk in r.iter_content(8192):
                f.write(chunk)
    return os.path.abspath(out_path)


def generate_media(
    media_type,
    text,
    session_id,
    api_key=None,
    template_path=None,
    aspect_ratio="16:9",
    style_type=None,
):
    """
    Generate different media types from text with comprehensive validation:
    - 'poster': Uses ideogram/ideogram-v3-turbo (Replicate) with auto-retry and DALL·E fallback
    - 'video': Uses bytedance/seedance-1-lite (Replicate)
    - 'slides': PPTX presentation (using OpenAI GPT)
    - 'memo': DOCX memo (using OpenAI GPT and python-docx)
    """
    # Validate media type
    if media_type not in ["poster", "video", "slides", "memo"]:
        raise ValueError(
            f"Invalid media type: {media_type}. Must be one of: poster, video, slides, memo"
        )

    # Validate aspect ratio
    if aspect_ratio not in VALID_RATIOS:
        raise ValueError(
            f"Invalid aspect ratio: {aspect_ratio}. Must be one of: {VALID_RATIOS}"
        )

    # Poster-specific validation
    if media_type == "poster":
        if template_path and style_type:
            raise ValueError(
                "Cannot specify both template_path and style_type - choose one"
            )

        # Set default style if neither provided
        if not template_path and not style_type:
            style_type = "DESIGN"

    timestamp = datetime.utcnow().strftime("%Y%m%d%H%M%S")
    safe_session = session_id or "anon"
    filename_base = f"{media_type}_{safe_session}_{timestamp}"

    # OpenAI API key required for slides and memo
    if media_type in ["slides", "memo"] and not api_key:
        raise RuntimeError("OpenAI API key is required for slides and memo generation")

    # Generate poster with auto-retry and fallback
    if media_type == "poster":
        if not REPLICATE_API_TOKEN:
            raise RuntimeError("REPLICATE_API_TOKEN environment variable not set")

        prompt = text.strip()[:400]  # Keep prompt short
        refs = None

        try:
            # Base input parameters
            input_data = {
                "prompt": prompt,
                "aspect_ratio": aspect_ratio,
            }

            # Process reference image if provided
            if template_path and os.path.exists(template_path):
                try:
                    refs = [open(template_path, "rb")]
                    input_data["style_reference_images"] = refs
                    logger.info(f"Using reference image: {template_path}")
                except FileNotFoundError:
                    logger.warning(f"Template file not found: {template_path}")
            # Process style type if provided
            elif style_type:
                valid_styles = ["AUTO", "GENERAL", "REALISTIC", "DESIGN"]
                if style_type.upper() not in valid_styles:
                    raise ValueError(
                        f"Invalid style type: {style_type}. Valid options: {valid_styles}"
                    )
                input_data["style_type"] = style_type.upper()
                logger.info(f"Using style type: {style_type.upper()}")

            # First try without version pin
            url = _run_replicate(POSTER_MODEL, input_data)
            return _download(url, session_id, "poster")

        except ReplicateError as e:
            # Handle version mismatch (422 status)
            if getattr(e, "status", None) == 422:
                try:
                    # Get latest model version
                    model = replicate.models.get(POSTER_MODEL.split(":")[0])
                    latest_version = model.versions.list()[0].id
                    pinned_model = f"{POSTER_MODEL.split(':')[0]}:{latest_version}"

                    # Retry with latest version
                    logger.info(f"Retrying with pinned version: {pinned_model}")
                    url = _run_replicate(pinned_model, input_data)
                    return _download(url, session_id, "poster")
                except Exception:
                    # Final fallback to DALL-E 3
                    logger.warning("Falling back to DALL-E 3")
                    url = _run_dalle(prompt)
                    return _download(url, session_id, "poster")
            else:
                raise RuntimeError(f"Poster generation failed: {e}") from e
        finally:
            if refs:
                for f in refs:
                    f.close()

    # Generate video using Seedance
    elif media_type == "video":
        if not REPLICATE_API_TOKEN:
            raise RuntimeError("REPLICATE_API_TOKEN environment variable not set")

        try:
            video_url = _run_replicate(
                "bytedance/seedance-1-lite",
                {
                    "prompt": text[:500],
                    "video_length": "5s",
                    "fps": 24,
                    "seed": 42,
                },
            )
            return _download(video_url, session_id, "video")
        except Exception as e:
            logger.error(f"Video generation error: {e}")
            raise RuntimeError(f"Video generation failed: {str(e)}")

    # Generate slides locally using OpenAI
    elif media_type == "slides":
        out_path = os.path.join(OUTPUT_FOLDER, f"{filename_base}.pptx")
        try:
            client = OpenAI(api_key=api_key)
            system_msg = {
                "role": "system",
                "content": "Convert this content into 3-5 slides in JSON format: "
                "[{'title':..., 'bullets':[...]}]",
            }
            resp = client.chat.completions.create(
                model="gpt-4",
                messages=[system_msg, {"role": "user", "content": text}],
                temperature=0.2,
            )
            slides = json.loads(resp.choices[0].message.content)
            prs = Presentation()

            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "Presentation"

            # Add content slides
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.get("title", "Slide")
                tf = slide.shapes.placeholders[1].text_frame
                for bullet in slide_data.get("bullets", []):
                    tf.add_paragraph().text = bullet

            prs.save(out_path)
            return os.path.abspath(out_path)
        except Exception as e:
            logger.error(f"Slide generation error: {e}")
            raise RuntimeError(f"Slide generation failed: {str(e)}")

    # Generate memo locally using OpenAI
    elif media_type == "memo":
        out_path = os.path.join(OUTPUT_FOLDER, f"{filename_base}.docx")
        try:
            client = OpenAI(api_key=api_key)
            resp = client.chat.completions.create(
                model="gpt-4",
                messages=[
                    {
                        "role": "user",
                        "content": f"Format this as a professional memo with clear sections:\n\n{text}",
                    }
                ],
                temperature=0.3,
            )
            memo_content = resp.choices[0].message.content
            doc = DocxWriter()

            # Add memo header
            doc.add_heading("MEMORANDUM", level=0)
            doc.add_paragraph().add_run(
                "Date: " + datetime.now().strftime("%B %d, %Y")
            ).bold = True
            doc.add_paragraph().add_run("To: ").bold = True
            doc.add_paragraph().add_run("From: ").bold = True
            doc.add_paragraph().add_run("Subject: ").bold = True
            doc.add_paragraph()

            # Add memo content
            for para in memo_content.splitlines():
                if para.strip():
                    doc.add_paragraph(para.strip())

            doc.save(out_path)
            return os.path.abspath(out_path)
        except Exception as e:
            logger.error(f"Memo generation error: {e}")
            raise RuntimeError(f"Memo generation failed: {str(e)}")


# ==============================================
# Text Extraction Functions
# ==============================================


def extract_text_from_zip(zip_path):
    """Extract text from zip archive"""
    extracted_text = ""
    with zipfile.ZipFile(zip_path, "r") as zip_ref:
        for file in zip_ref.namelist():
            if any(file.lower().endswith(ext) for ext in ALLOWED_EXTENSIONS):
                with zip_ref.open(file) as f:
                    try:
                        with tempfile.NamedTemporaryFile(delete=False) as tmp:
                            tmp.write(f.read())
                            tmp_path = tmp.name
                        extracted_text += extract_text_from_file(tmp_path) + "\n\n"
                        os.unlink(tmp_path)
                    except Exception as e:
                        logger.error(f"Error processing {file} in ZIP: {str(e)}")
    return extracted_text


def extract_text_from_file(filepath, api_key=None):
    """Extract text from various file types using GPT-4o vision when needed"""
    ext = Path(filepath).suffix.lower()
    text = ""

    try:
        if ext == ".pdf":
            if is_image_heavy_pdf(filepath):
                text = process_pdf_with_vision(filepath, api_key)
            else:
                text = extract_pdf_text(filepath)
        elif ext in [".jpg", ".jpeg", ".png"]:
            text = analyze_image_with_gpt4o(filepath, api_key)
        elif ext in [".doc", ".docx"]:
            text = extract_docx_text(filepath)
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()
    except Exception as e:
        logger.error(f"Error extracting text: {str(e)}")
        return ""
    return text


def process_pdf_with_vision(pdf_path, api_key):
    """Use GPT-4o vision to extract text from PDF pages"""
    client = OpenAI(api_key=api_key)
    doc = fitz.open(pdf_path)
    full_text = []

    for page_num in range(len(doc)):
        pix = doc[page_num].get_pixmap()
        img_path = f"page_{page_num}.png"
        pix.save(img_path)

        with open(img_path, "rb") as img_file:
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": "Extract all text from this document page, preserving layout, tables and structure.",
                    },
                    {"role": "user", "image": img_file.read()},
                ],
                max_tokens=2000,
            )
            full_text.append(response.choices[0].message.content)
        os.remove(img_path)
    return "\n\n".join(full_text)


def analyze_image_with_gpt4o(image_path, api_key):
    """Use GPT-4o vision to analyze images/charts"""
    client = OpenAI(api_key=api_key)
    with open(image_path, "rb") as img_file:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": "Describe this image in detail, including any text, charts, or diagrams.",
                },
                {"role": "user", "image": img_file.read()},
            ],
            max_tokens=1000,
        )
        return response.choices[0].message.content


# ==============================================
# NLP and AI Processing Functions
# ==============================================


def gpt_semantic_chunk(text, api_key, max_chunks=10):
    """Use GPT-4o to split text into coherent semantic chunks"""
    try:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": f"""
                    Split this document into up to {max_chunks} coherent semantic chunks.
                    Each chunk should represent a logically self-contained section or idea (~400-600 words).
                    Preserve all formatting, tables, and layout information.
                    Return a JSON list of chunks.
                    
                    Document:
                    {text[:6000]}  # Limit context window
                    """,
                }
            ],
            response_format={"type": "json_object"},
            temperature=0.3,
        )
        return json.loads(response.choices[0].message.content)["chunks"]
    except Exception as e:
        logger.error(f"GPT chunking failed: {str(e)}")
        return [text]  # Fallback to original text


def embed(text, api_key, model="text-embedding-3-large"):
    """Standard embedding function with fixed dimensions"""
    client = OpenAI(api_key=api_key)
    response = client.embeddings.create(
        model=model,
        input=[text],
        encoding_format="float",
        dimensions=3072,  # Explicitly set dimensions
    )
    return response.data[0].embedding


def get_enhanced_embeddings(texts, api_key, model="text-embedding-3-large"):
    """Generate embeddings with cross-chunk context"""
    if not isinstance(texts, list):
        texts = [texts]

    try:
        client = OpenAI(api_key=api_key)
        summary_prompt = f"""
        Analyze these text chunks as a cohesive document and create a global summary 
        that captures the main themes and relationships between chunks:
        {texts}
        """
        summary = (
            client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": summary_prompt}],
                max_tokens=500,
            )
            .choices[0]
            .message.content
        )

        embeddings = []
        for text in texts:
            enhanced_text = f"Chunk: {text}\n\nGlobal Context: {summary}"
            embedding = embed(enhanced_text, api_key, model)
            embeddings.append(embedding)

        return embeddings[0] if len(embeddings) == 1 else embeddings
    except Exception as e:
        logger.error(f"Enhanced embeddings failed: {str(e)}")
        return embed(texts, api_key, model)  # Fallback to regular embedding


def detect_document_type(text, api_key):
    """Use GPT to classify document type"""
    try:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": f"""
                    Classify this document into one of these types:
                    [memo, article, report, slide, code, legal, spreadsheet, general]
                    
                    Text:
                    {text[:2000]}
                    
                    Return only the type.
                    """,
                }
            ],
            max_tokens=10,
            temperature=0,
        )
        return response.choices[0].message.content.strip().lower()
    except Exception as e:
        logger.error(f"Document type detection failed: {str(e)}")
        return "general"


def extract_concepts(text, api_key):
    client = OpenAI(api_key=api_key)
    prompt = (
        "Extract key concepts and summarize the main ideas from the following text:\n"
        + text[:4000]
    )
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300,
    )
    return response.choices[0].message.content


def perform_web_search(query, api_key):
    """Simulate or perform actual web search"""
    client = OpenAI(api_key=api_key)
    prompt = f"""
    You are a web search assistant. Given the query "{query}", generate 3-5 realistic 
    search results with titles, URLs, and snippets as if they came from a real search engine.
    
    Return the results in JSON format like this:
    {{
        "results": [
            {{
                "title": "Result title",
                "source": "example.com",
                "snippet": "Relevant information snippet..."
            }}
        ]
    }}
    """
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
    )
    data = json.loads(response.choices[0].message.content)
    return data.get("results", [])


def get_preview(media_type, prompt, api_key):
    client = OpenAI(api_key=api_key)
    if media_type == "image":
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "user", "content": f"Describe a thumbnail image for: {prompt}"}
            ],
            max_tokens=150,
        )
        return response.choices[0].message.content
    elif media_type == "slides":
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "user",
                    "content": f"Create an outline for slides based on: {prompt}",
                }
            ],
            max_tokens=300,
        )
        return response.choices[0].message.content
    elif media_type == "video":
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "user",
                    "content": f"Describe a short video preview for: {prompt}",
                }
            ],
            max_tokens=200,
        )
        return response.choices[0].message.content
    return "Preview not available for this media type"


def initialize_vector_index(path="vector_index"):
    """Initialize or load the FAISS index"""
    index_path = f"{path}.faiss"
    meta_path = f"{path}.json"

    if os.path.exists(index_path):
        index = faiss.read_index(index_path)
        with open(meta_path, "r") as f:
            metadata = json.load(f)
    else:
        index = faiss.IndexIDMap(faiss.IndexFlatL2(3072))
        metadata = {}

    return index, metadata


def save_vector_index(index, metadata, path="vector_index"):
    """Save the FAISS index"""
    index_path = f"{path}.faiss"
    meta_path = f"{path}.json"

    faiss.write_index(index, index_path)
    with open(meta_path, "w") as f:
        json.dump(metadata, f, indent=2)

    increment_index_version()


def add_to_index(index, metadata, vector, vector_metadata):
    """Add a vector to the index with metadata"""
    vector_id = len(metadata) + 1
    index.add_with_ids(
        np.array([vector]).astype("float32"), np.array([vector_id], dtype=np.int64)
    )
    metadata[vector_id] = vector_metadata
    return vector_id


def remove_from_index(index, metadata, vector_id):
    """Remove a vector from the index"""
    index.remove_ids(np.array([vector_id], dtype=np.int64))
    if vector_id in metadata:
        del metadata[vector_id]


def search_index(index, metadata, query_vector, k=5):
    """Search the index for similar vectors"""
    distances, vector_ids = index.search(np.array([query_vector]).astype("float32"), k)
    results = []
    for i in range(len(vector_ids[0])):
        if vector_ids[0][i] >= 0 and vector_ids[0][i] in metadata:
            results.append(
                {**metadata[vector_ids[0][i]], "similarity": 1 - distances[0][i]}
            )
    return results


# ======================
# RAG Pipeline Functions
# ======================


def generate_answer(question, context, api_key):
    """Generate an answer using GPT with provided context"""
    client = OpenAI(api_key=api_key)

    # Format context for the prompt
    context_str = "\n\n".join(
        f"[Source {i+1}]: {res['content']}" for i, res in enumerate(context)
    )

    messages = [
        {
            "role": "system",
            "content": "You are a knowledgeable assistant. Use the provided context to answer questions accurately and concisely.",
        },
        {"role": "user", "content": f"Context:\n{context_str}\n\nQuestion: {question}"},
    ]

    try:
        response = client.chat.completions.create(
            model=GPT_MODEL, messages=messages, temperature=0.3, max_tokens=1024
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"Answer generation failed: {str(e)}")
        return "Error generating answer"


def main():
    """Main RAG pipeline execution"""
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        logger.error("OPENAI_API_KEY environment variable not set")
        return

    # Get question from environment or use default
    question = os.getenv("QUESTION", "What is the main topic of these documents?")

    # Check for existing index or create new one
    index_path = Path(OUTPUT_FOLDER) / "index.faiss"
    metadata_path = Path(OUTPUT_FOLDER) / "metadata.json"

    if index_path.exists() and metadata_path.exists():
        logger.info("Loading existing index")
        index = faiss.read_index(str(index_path))
        with open(metadata_path, "r") as f:
            metadata = json.load(f)
    else:
        logger.info("Building new index")
        index, metadata = create_index(api_key)
        faiss.write_index(index, str(index_path))
        with open(metadata_path, "w") as f:
            json.dump(metadata, f)

    # Process the question
    logger.info("Processing question")
    question_embedding = get_embeddings([question], api_key)[0]
    context = search_index(index, metadata, question_embedding)
    answer = generate_answer(question, context, api_key)

    # Save results
    result = {
        "question": question,
        "answer": answer,
        "sources": [
            {"file": item["filename"], "chunk": item["chunk_index"]} for item in context
        ],
        "timestamp": datetime.now().isoformat(),
    }

    with open(Path(OUTPUT_FOLDER) / "result.json", "w") as f:
        json.dump(result, f, indent=2)

    logger.info(f"Answer: {answer}")
    logger.info("RAG pipeline completed")


if __name__ == "__main__":
    main()
