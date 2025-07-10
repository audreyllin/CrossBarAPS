from flask import Flask, request, jsonify, render_template, send_file
import os
import logging
import uuid
import tiktoken
import json
import PyPDF2
import pdfplumber
from docx import Document as DocxDoc
import hashlib
import re
import sqlite3
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
from openai import OpenAI
from werkzeug.utils import secure_filename
from collections import Counter
import pytesseract
from PIL import Image
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from io import BytesIO
import datetime

# Initialize app and paths
current_dir = os.path.dirname(os.path.abspath(__file__))
template_path = os.path.join(current_dir, "templates")
DB_FILE = os.path.join(current_dir, "memory.db")
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
ALLOWED_EXTENSIONS = {
    "txt",
    "pdf",
    "doc",
    "docx",
    "ppt",
    "pptx",
    "jpg",
    "jpeg",
    "png",
    "mp4",
}

# Admin credentials
ADMIN_USER = os.environ.get("ADMIN_USER", "admin")
ADMIN_PASS = os.environ.get("ADMIN_PASS", "securepassword")

app = Flask(__name__, template_folder=template_path)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
logging.basicConfig(level=logging.INFO)


# Initialize database
def init_db():
    with sqlite3.connect(DB_FILE) as db:
        db.execute(
            """
            CREATE TABLE IF NOT EXISTS context_vectors (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                context_id TEXT NOT NULL,
                session_id TEXT NOT NULL,
                vector TEXT NOT NULL,
                model TEXT NOT NULL,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        db.execute(
            """
            CREATE TABLE IF NOT EXISTS session_contexts (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT NOT NULL,
                context_id TEXT NOT NULL,
                content TEXT NOT NULL,
                content_hash TEXT NOT NULL,
                source_type TEXT,
                filename TEXT,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP,
                is_active BOOLEAN DEFAULT 1
            )
            """
        )
        db.execute(
            """
            CREATE TABLE IF NOT EXISTS conversation_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                session_id TEXT NOT NULL,
                question TEXT NOT NULL,
                answer TEXT NOT NULL,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
            )
            """
        )
        # Add columns if missing
        try:
            db.execute("ALTER TABLE session_contexts ADD COLUMN content_hash TEXT")
        except sqlite3.OperationalError:
            pass
        try:
            db.execute(
                "ALTER TABLE session_contexts ADD COLUMN is_active BOOLEAN DEFAULT 1"
            )
        except sqlite3.OperationalError:
            pass
        # Create indexes
        db.execute(
            "CREATE INDEX IF NOT EXISTS idx_content_hash ON session_contexts (content_hash)"
        )
        db.execute(
            "CREATE INDEX IF NOT EXISTS idx_session_active ON session_contexts (session_id, is_active)"
        )


init_db()

# Token counting
encoding = tiktoken.encoding_for_model("gpt-3.5-turbo")


def count_tokens(text):
    return len(encoding.encode(text))


# Text extraction from files
def extract_text_from_file(filepath, filename):
    ext = filename.split(".")[-1].lower()
    text = ""
    try:
        if ext == "pdf":
            with pdfplumber.open(filepath) as pdf:
                text = "\n".join(
                    [page.extract_text() for page in pdf.pages if page.extract_text()]
                )
        elif ext in ["doc", "docx"]:
            doc = Document(filepath)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif ext in ["ppt", "pptx"]:
            prs = Presentation(filepath)
            text = "\n".join(
                [
                    shape.text
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                ]
            )
        elif ext in ["jpg", "jpeg", "png"]:
            image = Image.open(filepath)
            text = pytesseract.image_to_string(image)
        elif ext == "mp4":
            text = "Video content not extracted"
        else:  # txt and others
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()
    except Exception as e:
        logging.error(f"Text extraction failed: {str(e)}")
        text = f"Error extracting text: {str(e)}"
    return text[:20000]  # Limit to 20k characters


# GPT Embedding
def generate_vector(text, api_key, model="text-embedding-3-large"):
    try:
        client = OpenAI(api_key=api_key)
        response = client.embeddings.create(
            model=model, input=[text[:8192]]  # Truncate to model limit
        )
        return np.array(response.data[0].embedding)
    except Exception as e:
        logging.error(f"Embedding generation failed: {str(e)}")
        dim = 3072 if model == "text-embedding-3-large" else 1536
        return np.zeros(dim)


# Concept extraction
def extract_concepts(text, api_key):
    try:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {
                    "role": "system",
                    "content": "Extract top 5 key concepts as JSON array",
                },
                {"role": "user", "content": f"Extract concepts from:\n\n{text[:3000]}"},
            ],
            max_tokens=150,
            temperature=0.3,
        )
        concepts = json.loads(response.choices[0].message.content)
        if isinstance(concepts, list) and len(concepts) > 0:
            return [concept.capitalize() for concept in concepts]
    except Exception as e:
        logging.error(f"Concept extraction failed: {str(e)}")

    # Fallback method
    words = re.findall(r"\b\w{4,}\b", text.lower())
    stop_words = {
        "this",
        "that",
        "which",
        "with",
        "from",
        "your",
        "have",
        "more",
        "about",
    }
    top_words = [
        word
        for word, count in Counter(words).most_common(5)
        if word not in stop_words and len(word) > 3
    ]
    return [word.capitalize() for word in top_words]


# Context retrieval
def get_context_vectors(session_id, model):
    with sqlite3.connect(DB_FILE) as db:
        rows = db.execute(
            "SELECT context_id, vector FROM context_vectors WHERE session_id = ? AND model = ?",
            (session_id, model),
        ).fetchall()
    return [(row[0], np.array(json.loads(row[1]))) for row in rows]


def get_relevant_contexts(
    question, session_contexts, session_id, api_key, model, top_n=3
):
    if not question or not session_contexts:
        return []

    stored_vectors = get_context_vectors(session_id, model)
    if not stored_vectors:
        return []

    question_vector = generate_vector(question, api_key, model=model)
    similarities = []
    for context_id, ctx_vec in stored_vectors:
        try:
            similarity = cosine_similarity([question_vector], [ctx_vec])[0][0]
            similarities.append((context_id, similarity))
        except Exception as e:
            logging.error(f"Similarity error: {str(e)}")
            similarities.append((context_id, 0))

    similarities.sort(key=lambda x: x[1], reverse=True)
    top_contexts = [ctx_id for ctx_id, _ in similarities[:top_n]]
    return [ctx for ctx in session_contexts if ctx["id"] in top_contexts]


# ===== ADMIN ROUTES =====
@app.route("/api/admin_login", methods=["POST"])
def admin_login():
    data = request.get_json()
    if data.get("username") == ADMIN_USER and data.get("password") == ADMIN_PASS:
        return jsonify({"status": "success"})
    return jsonify({"error": "Invalid credentials"}), 401


@app.route("/api/admin/reset_db", methods=["POST"])
def reset_db():
    try:
        os.system("python init_db.py")
        return jsonify({"status": "reset successful"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/admin/conversations", methods=["GET"])
def admin_conversations():
    with sqlite3.connect(DB_FILE) as db:
        rows = db.execute(
            "SELECT session_id, question, answer, timestamp FROM conversation_history ORDER BY timestamp DESC"
        ).fetchall()
    return jsonify(
        [
            {"session_id": r[0], "question": r[1], "answer": r[2], "timestamp": r[3]}
            for r in rows
        ]
    )


@app.route("/api/admin/frequent_questions", methods=["GET"])
def frequent_questions():
    with sqlite3.connect(DB_FILE) as db:
        rows = db.execute(
            "SELECT question, COUNT(*) as count FROM conversation_history GROUP BY question ORDER BY count DESC LIMIT 10"
        ).fetchall()
    return jsonify([{"question": r[0], "count": r[1]} for r in rows])


# ===== FILE MANAGEMENT ROUTES =====
@app.route("/api/download_conversation", methods=["GET"])
def download_convo():
    session_id = request.args.get("sessionId")
    fmt = request.args.get("format", "txt")
    if not session_id:
        return jsonify({"error": "Missing session ID"}), 400

    with sqlite3.connect(DB_FILE) as db:
        rows = db.execute(
            "SELECT question, answer, timestamp FROM conversation_history WHERE session_id = ?",
            (session_id,),
        ).fetchall()

    content = [f"Q: {q}\nA: {a}\nTime: {t}" for q, a, t in rows]

    if fmt == "txt":
        return (
            BytesIO("\n\n".join(content).encode("utf-8")),
            200,
            {
                "Content-Type": "text/plain",
                "Content-Disposition": f"attachment; filename=chat_{session_id}.txt",
            },
        )
    elif fmt == "docx":
        doc = DocxDoc()
        for entry in content:
            doc.add_paragraph(entry)
        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)
        return send_file(
            buf, as_attachment=True, download_name=f"chat_{session_id}.docx"
        )
    elif fmt == "pdf":
        buf = BytesIO()
        c = canvas.Canvas(buf, pagesize=letter)
        width, height = letter
        y = height - 50
        for entry in content:
            for line in entry.split("\n"):
                c.drawString(50, y, line)
                y -= 15
                if y < 50:
                    c.showPage()
                    y = height - 50
        c.save()
        buf.seek(0)
        return send_file(
            buf, as_attachment=True, download_name=f"chat_{session_id}.pdf"
        )
    return jsonify({"error": "Unsupported format"}), 400


@app.route("/api/list_contexts", methods=["GET"])
def list_contexts():
    session_id = request.args.get("sessionId")
    if not session_id:
        return jsonify({"error": "Missing session ID"}), 400

    with sqlite3.connect(DB_FILE) as db:
        rows = db.execute(
            "SELECT context_id, filename, timestamp FROM session_contexts "
            "WHERE session_id = ? AND is_active = 1",
            (session_id,),
        ).fetchall()
    return jsonify(
        [{"context_id": r[0], "filename": r[1], "timestamp": r[2]} for r in rows]
    )


@app.route("/api/remove_context", methods=["POST"])
def remove_context():
    data = request.get_json()
    context_id = data.get("contextId")
    session_id = data.get("sessionId")
    if not context_id or not session_id:
        return jsonify({"error": "Missing parameters"}), 400

    with sqlite3.connect(DB_FILE) as db:
        db.execute(
            "UPDATE session_contexts SET is_active = 0 WHERE context_id = ? AND session_id = ?",
            (context_id, session_id),
        )
        db.execute(
            "DELETE FROM context_vectors WHERE context_id = ? AND session_id = ?",
            (context_id, session_id),
        )
        db.commit()
    return jsonify({"status": "removed"})


# ===== MAIN FUNCTIONALITY ROUTES =====
@app.route("/api/upload_context", methods=["POST"])
def upload_context():
    api_key = request.headers.get("Authorization", "").replace("Bearer ", "")
    if not api_key:
        return jsonify({"error": "Missing API key"}), 401

    session_id = request.form.get("sessionId")
    model = request.form.get("model", "text-embedding-3-large")
    results = []

    # Handle text context
    if "text" in request.form:
        text_content = request.form["text"].strip()[:20000]
        if not text_content:
            return jsonify({"error": "Empty text content"}), 400

        content_hash = hashlib.sha256(text_content.encode()).hexdigest()

        # Check for duplicates
        with sqlite3.connect(DB_FILE) as db:
            existing = db.execute(
                "SELECT context_id FROM session_contexts "
                "WHERE session_id = ? AND content_hash = ? AND is_active = 1",
                (session_id, content_hash),
            ).fetchone()

            if existing:
                return jsonify(
                    {
                        "status": "duplicate",
                        "message": "Content already exists",
                        "context_id": existing[0],
                    }
                )

        context_id = f"text_{hashlib.md5(text_content.encode()).hexdigest()[:10]}"
        vector = generate_vector(text_content, api_key, model=model)

        with sqlite3.connect(DB_FILE) as db:
            db.execute(
                "INSERT INTO context_vectors (context_id, session_id, vector, model) "
                "VALUES (?, ?, ?, ?)",
                (context_id, session_id, json.dumps(vector.tolist()), model),
            )
            db.execute(
                "INSERT INTO session_contexts (session_id, context_id, content, source_type, content_hash) "
                "VALUES (?, ?, ?, ?, ?)",
                (session_id, context_id, text_content, "text", content_hash),
            )
            db.commit()

        concepts = extract_concepts(text_content, api_key)
        results.append(
            {"status": "success", "context_id": context_id, "concepts": concepts}
        )

    # Handle file context
    elif "file" in request.files:
        files = request.files.getlist("file")
        for file in files:
            if file.filename == "":
                continue

            if not allowed_file(file.filename):
                continue

            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)
            text_content = extract_text_from_file(filepath, filename)
            content_hash = hashlib.sha256(text_content.encode()).hexdigest()

            # Check for duplicates
            with sqlite3.connect(DB_FILE) as db:
                existing = db.execute(
                    "SELECT context_id FROM session_contexts "
                    "WHERE session_id = ? AND content_hash = ? AND is_active = 1",
                    (session_id, content_hash),
                ).fetchone()

                if existing:
                    os.remove(filepath)
                    results.append(
                        {
                            "status": "duplicate",
                            "filename": filename,
                            "context_id": existing[0],
                        }
                    )
                    continue

            context_id = (
                f"file_{hashlib.md5((filename + session_id).encode()).hexdigest()[:10]}"
            )
            vector = generate_vector(text_content, api_key, model=model)

            with sqlite3.connect(DB_FILE) as db:
                db.execute(
                    "INSERT INTO context_vectors (context_id, session_id, vector, model) "
                    "VALUES (?, ?, ?, ?)",
                    (context_id, session_id, json.dumps(vector.tolist()), model),
                )
                db.execute(
                    "INSERT INTO session_contexts (session_id, context_id, content, source_type, filename, content_hash) "
                    "VALUES (?, ?, ?, ?, ?, ?)",
                    (
                        session_id,
                        context_id,
                        text_content,
                        "file",
                        filename,
                        content_hash,
                    ),
                )
                db.commit()

            concepts = extract_concepts(text_content, api_key)
            results.append(
                {"filename": filename, "context_id": context_id, "concepts": concepts}
            )
    else:
        return jsonify({"error": "No valid context provided"}), 400

    return jsonify(results)


@app.route("/api/ask", methods=["POST"])
def ask_question():
    api_key = request.headers.get("Authorization", "").replace("Bearer ", "")
    if not api_key:
        return jsonify({"error": "Missing API key"}), 401

    data = request.get_json()
    question = data.get("question", "").strip()
    session_id = data.get("sessionId", "")
    model = data.get("model", "text-embedding-3-large")

    if not question:
        return jsonify({"error": "Missing question"}), 400
    if not session_id:
        return jsonify({"error": "Missing session ID"}), 400

    try:
        # Get active contexts
        with sqlite3.connect(DB_FILE) as db:
            rows = db.execute(
                "SELECT context_id, content FROM session_contexts "
                "WHERE session_id = ? AND is_active = 1",
                (session_id,),
            ).fetchall()
        session_contexts = [{"id": row[0], "content": row[1]} for row in rows]

        # Get relevant contexts
        relevant_contexts = get_relevant_contexts(
            question, session_contexts, session_id, api_key, model
        )

        # Prepare context for GPT
        context_text = "\n\n".join([ctx["content"] for ctx in relevant_contexts[:3]])
        prompt = f"Context:\n{context_text}\n\nQuestion: {question}\nAnswer:"

        # Call GPT
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "Answer based on provided context"},
                {"role": "user", "content": prompt},
            ],
            max_tokens=500,
            temperature=0.3,
        )
        answer = response.choices[0].message.content.strip()

        # Store conversation
        with sqlite3.connect(DB_FILE) as db:
            db.execute(
                "INSERT INTO conversation_history (session_id, question, answer) "
                "VALUES (?, ?, ?)",
                (session_id, question, answer),
            )
            db.commit()

        # Generate follow-up questions
        follow_up_prompt = f"Given Q: '{question}' and A: '{answer}', generate 3 relevant follow-ups as JSON array."
        follow_up_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": follow_up_prompt}],
            max_tokens=150,
            temperature=0.5,
        )

        try:
            follow_ups = json.loads(follow_up_response.choices[0].message.content)
        except:
            follow_ups = [
                "Can you elaborate?",
                "Related topics?",
                "What else should I know?",
            ]

        return jsonify(
            {
                "question": question,
                "answer": answer,
                "contexts": [ctx["content"] for ctx in relevant_contexts],
                "follow_ups": follow_ups,
            }
        )
    except Exception as e:
        logging.error(f"Question answering failed: {str(e)}")
        return jsonify({"error": str(e)}), 500


# Helper functions
def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


@app.route("/")
def index():
    return render_template("kaggleloc.html")


def process_existing_uploads(
    session_id="manual_batch_session", model="text-embedding-3-large"
):
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        logging.warning("OPENAI_API_KEY not set. Skipping file auto-processing.")
        return

    processed = 0
    skipped = 0

    for fname in os.listdir(UPLOAD_FOLDER):
        if not allowed_file(fname):
            continue
        fpath = os.path.join(UPLOAD_FOLDER, fname)

        try:
            content = extract_text_from_file(fpath, fname)[:20000]
            content_hash = hashlib.sha256(content.encode()).hexdigest()

            # Skip if already in DB
            with sqlite3.connect(DB_FILE) as db:
                exists = db.execute(
                    "SELECT 1 FROM session_contexts WHERE session_id = ? AND content_hash = ? AND is_active = 1",
                    (session_id, content_hash),
                ).fetchone()
            if exists:
                logging.info(f"Skipped (already processed): {fname}")
                skipped += 1
                continue

            context_id = (
                f"file_{hashlib.md5((fname + session_id).encode()).hexdigest()[:10]}"
            )
            vector = generate_vector(content, api_key, model)
            concepts = extract_concepts(content, api_key)

            with sqlite3.connect(DB_FILE) as db:
                db.execute(
                    "INSERT INTO context_vectors (context_id, session_id, vector, model) VALUES (?, ?, ?, ?)",
                    (context_id, session_id, json.dumps(vector.tolist()), model),
                )
                db.execute(
                    "INSERT INTO session_contexts (session_id, context_id, content, source_type, filename, content_hash, is_active) VALUES (?, ?, ?, 'file', ?, ?, 1)",
                    (session_id, context_id, content, fname, content_hash),
                )
                db.commit()

            logging.info(f"Processed: {fname} | Concepts: {concepts}")
            processed += 1

        except Exception as e:
            logging.error(f"Failed on {fname}: {e}")

    logging.info(
        f"File auto-processing complete. Processed: {processed}, Skipped: {skipped}"
    )


if __name__ == "__main__":
    process_existing_uploads()  # Auto-process files on startup
    app.run(host="0.0.0.0", port=5000, debug=True)
